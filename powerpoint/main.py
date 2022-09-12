import os
if not os.path.exists("output"):
    os.mkdir("output")

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData, XyChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import numpy as np
import matplotlib.pyplot as plt

def title(prs):
    # Get a reference to a pre-existing slide layout
    title_slide_layout = prs.slide_layouts[TITLE]
    # Create a slide using the given layout
    slide = prs.slides.add_slide(title_slide_layout)
    # Get components for updating text
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # Update text
    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

def titleAndContent(prs):
    bullet_slide_layout = prs.slide_layouts[TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Adding a Bullet Slide'
    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'
    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    # Indent the paragraph
    p.level = 1
    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2

def custom(prs):
    blank_slide_layout = prs.slide_layouts[BLANK]
    slide = prs.slides.add_slide(blank_slide_layout)
    # Use PowerPoint position measurements
    margin = Inches(1)
    # Adds a textbox with 1in margin on all sides
    txBox = slide.shapes.add_textbox(margin, margin, prs.slide_width - margin * 2, prs.slide_height - margin * 2)
    tf = txBox.text_frame
    tf.text = "This is text inside a textbox"
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    # Use PowerPoint font measurements
    p.font.size = Pt(40)

def images(prs):
    img_path = 'hic.jpeg'
    blank_slide_layout = prs.slide_layouts[BLANK]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)
    left = Inches(4)
    height = Inches(4)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

def shapes(prs):
    title_only_slide_layout = prs.slide_layouts[TITLE_ONLY]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'Adding an AutoShape'
    # Add step 1
    left = Inches(0.93)
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)
    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    left = left + width - Inches(0.4)
    width = Inches(2.0) 
    # Add steps 2 through 6
    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        left = left + width - Inches(0.4)

def table(prs):
    title_only_slide_layout = prs.slide_layouts[TITLE_ONLY]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'Adding a Table'
    rows = cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)
    # write column headings
    table.cell(0, 0).text = 'Foo'
    table.cell(0, 1).text = 'Bar'
    # write body cells
    table.cell(1, 0).text = 'Baz'
    table.cell(1, 1).text = 'Qux'

def bar(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    # Prepare the chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
    # Add the chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    return slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

def cluster(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    chart_data = ChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
    chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
    chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    return slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

def scatter(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    chart_data = XyChartData()
    series_1 = chart_data.add_series('Model 1')
    series_1.add_data_point(0.7, 2.7)
    series_1.add_data_point(1.8, 3.2)
    series_1.add_data_point(2.6, 0.8)
    series_2 = chart_data.add_series('Model 2')
    series_2.add_data_point(1.3, 3.7)
    series_2.add_data_point(2.7, 2.3)
    series_2.add_data_point(1.6, 1.8)
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    )

def bubble(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    chart_data = BubbleChartData()
    series_1 = chart_data.add_series('Series 1')
    series_1.add_data_point(0.7, 2.7, 10)
    series_1.add_data_point(1.8, 3.2, 4)
    series_1.add_data_point(2.6, 0.8, 8)
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
    )

def axes(prs):
    chart = bar(prs)# Style the category axis
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = True
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(24)
    # Style the value axis
    value_axis = chart.value_axis
    value_axis.maximum_scale = 50.0
    value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_minor_gridlines = True
    # Style the tick labels
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(14)

def labels(prs):
    chart = bar(prs)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(13)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

def legend(prs):
    chart = cluster(prs)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

def line(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    chart_data = ChartData()
    chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
    chart_data.add_series('West',    (32.2, 28.4, 34.7))
    chart_data.add_series('East',    (24.3, 30.6, 20.2))
    chart_data.add_series('Midwest', (20.4, 18.3, 26.2))
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True

def pie(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    chart_data = ChartData()
    chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
    chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

def plot3d(prs):
    fig = plt.figure()
    ax = plt.axes(projection='3d')
    # Data for a three-dimensional line
    zline = np.linspace(0, 15, 1000)
    xline = np.sin(zline)
    yline = np.cos(zline)
    ax.plot3D(xline, yline, zline, 'gray')
    # Data for three-dimensional scattered points
    zdata = 15 * np.random.random(100)
    xdata = np.sin(zdata) + 0.1 * np.random.randn(100)
    ydata = np.cos(zdata) + 0.1 * np.random.randn(100)
    ax.scatter3D(xdata, ydata, zdata, c=zdata, cmap='Greens')
    # save
    img_path = 'output/3d.png'
    plt.savefig(img_path)

    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    left = top = Inches(1)
    width = prs.slide_width - Inches(2)
    height = prs.slide_height - Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)

TITLE = 0
TITLE_AND_CONTENT = 1
SECTION_HEADER = 2
TWO_CONTENT = 3
COMPARISON = 4
TITLE_ONLY = 5
BLANK = 6
CONTENT_WITH_CAPTION = 7
PICTURE_WITH_CAPTION = 8
TITLE_AND_VERTICAL_TEXT = 9
VERTICAL_TITLE_AND_TEXT = 10

if __name__ == "__main__":
    prs = Presentation()

    title(prs)
    titleAndContent(prs)
    custom(prs)
    images(prs)
    shapes(prs)
    table(prs)
    # charts
    bar(prs)
    cluster(prs)
    scatter(prs)
    bubble(prs)
    axes(prs)
    labels(prs)
    legend(prs)
    line(prs)
    pie(prs)
    plot3d(prs)
    for layout in prs.slide_layouts:
        slide = prs.slides.add_slide(layout)
        txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, Inches(1))
        tf = txBox.text_frame
        tf.text = layout.name
        tf.size = Pt(12)

    # save
    prs.save('output/test.pptx')